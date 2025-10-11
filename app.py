import streamlit as st
import requests
import pandas as pd
import os
import plotly.express as px
import matplotlib.pyplot as plt
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches
import tempfile
import plotly.io as pio

st.set_page_config(page_title="Social Media Analysis Master App", layout="wide")
st.title("Social Media Analysis Master App")

# TMU logo display
logo_path = os.path.join(os.path.dirname(__file__), "tmu_logo.png")
if os.path.exists(logo_path):
    st.image(logo_path, width=200)
else:
    st.write("TMU Logo not found")


def save_plotly_fig(fig, filename):
    pio.write_image(fig, filename)


def create_pptx_report(df, logo_path, output_path):
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "TMU Social Media Analysis Report"
    slide.placeholders[1].text = "Generated October 2025"
    if os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, Inches(0.5), Inches(0.5), width=Inches(1.5))
    with tempfile.TemporaryDirectory() as tmpdirname:
        fig1 = px.line(df, x='timestamp', y='engagement_rate', title='Engagement Rate Over Time')
        path1 = os.path.join(tmpdirname, "engagement_rate.png")
        save_plotly_fig(fig1, path1)
        metrics = ['likes', 'comments', 'shares']
        sums = df[metrics].sum().reset_index()
        sums.columns = ['Metric', 'Count']
        fig2 = px.bar(sums, x='Metric', y='Count', title='Total Likes, Comments & Shares')
        path2 = os.path.join(tmpdirname, "likes_comments_shares.png")
        save_plotly_fig(fig2, path2)
        for title, img_path in {"Engagement Rate Over Time": path1, "Likes, Comments & Shares": path2}.items():
            slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = title
            slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), width=Inches(7), height=Inches(4))
        prs.save(output_path)


tabs = st.tabs([
    "Upload", "Overview", "Preview", "Dashboard", "Audience",
    "Content Insights", "Best Time", "Advanced Analytics", "Hashtag Analysis",
    "Performance Trends", "ML Explainability", "Forecast", "Clustering", "Export"
])

# Upload Tab
with tabs[0]:
    st.header("Upload Your Social Media CSV Files")
    uploaded_files = st.file_uploader(
        "Select one or more CSV files to upload",
        type=["csv"],
        accept_multiple_files=True
    )

    if uploaded_files:
        for uploaded_file in uploaded_files:
            response = requests.post(
                "http://backend:8000/upload/",
                files={"file": (uploaded_file.name, uploaded_file, "text/csv")}
            )
            result = response.json()
            if "rows" in result:
                st.success(f"Successfully uploaded {uploaded_file.name}: {result['rows']} rows")
            else:
                st.error(f"Failed to upload {uploaded_file.name}: {result.get('error', 'Unknown error')}")

# Overview Tab
with tabs[1]:
    st.header("Analysis Overview")
    response = requests.get("http://backend:8000/analysis/overview/")
    if response.status_code == 200:
        overview = response.json()
        if "error" not in overview:
            col1, col2, col3, col4 = st.columns(4)
            with col1:
                st.metric("Total Posts", overview['total_posts'])
            with col2:
                st.metric("Avg Engagement Rate", f"{overview['avg_engagement_rate']:.3f}")
            with col3:
                st.metric("Total Likes", overview['total_likes'])
            with col4:
                st.metric("Total Comments", overview['total_comments'])

            col1, col2, col3 = st.columns(3)
            with col1:
                st.metric("Total Shares", overview['total_shares'])
            with col2:
                st.metric("Best Post Type", overview['best_performing_post_type'])
            with col3:
                st.metric("Peak Hour", f"{overview['peak_engagement_hour']}:00")

# Preview Tab
with tabs[2]:
    st.header("Data Preview")
    if st.button("Load Preview"):
        preview = requests.get("http://backend:8000/preview/")
        df = pd.DataFrame(preview.json())
        st.dataframe(df)

# Dashboard Tab
with tabs[3]:
    st.header("Engagement Dashboard")
    preview = requests.get("http://backend:8000/preview/")
    df = pd.DataFrame(preview.json())
    if len(df) > 0:
        df['timestamp'] = pd.to_datetime(df['timestamp'])
        col1, col2 = st.columns(2)
        with col1:
            st.subheader("Engagement Rate Over Time")
            fig = px.line(df, x='timestamp', y='engagement_rate', title='Engagement Rate Over Time')
            st.plotly_chart(fig, use_container_width=True)
        with col2:
            st.subheader("Performance Score Distribution")
            fig2 = px.histogram(df, x='performance_score', title='Performance Score Distribution')
            st.plotly_chart(fig2, use_container_width=True)

# Advanced Analytics Tab
with tabs[7]:
    st.header("Advanced Analytics")

    # Audience Segments
    st.subheader("Audience Segments")
    response = requests.get("http://backend:8000/analysis/segments/")
    if response.status_code == 200:
        segments = response.json()
        df_segments = pd.DataFrame(segments)
        fig = px.bar(df_segments, x='audience_segment', y='count', title='Posts by Audience Segment')
        st.plotly_chart(fig, use_container_width=True)

    # Content Length Analysis
    st.subheader("Content Length Performance")
    response = requests.get("http://backend:8000/analysis/content-insights/")
    if response.status_code == 200:
        content_insights = response.json()
        df_content = pd.DataFrame(content_insights)
        fig = px.bar(df_content, x='content_length_category', y='avg_engagement',
                     title='Average Engagement by Content Length')
        st.plotly_chart(fig, use_container_width=True)

# Hashtag Analysis Tab
with tabs[8]:
    st.header("Hashtag Analysis")
    response = requests.get("http://backend:8000/analysis/top-hashtags/")
    if response.status_code == 200:
        hashtags = response.json()
        if hashtags:
            df_hashtags = pd.DataFrame(list(hashtags.items()), columns=['Hashtag', 'Count'])
            fig = px.bar(df_hashtags, x='Hashtag', y='Count', title='Top 10 Hashtags')
            st.plotly_chart(fig, use_container_width=True)

# Performance Trends Tab
with tabs[9]:
    st.header("Performance Trends")
    response = requests.get("http://backend:8000/analysis/performance-trends/")
    if response.status_code == 200:
        trends = response.json()
        if trends:
            df_trends = pd.DataFrame(trends)
            df_trends['date'] = pd.to_datetime(df_trends['date'])
            fig = px.line(df_trends, x='date', y='avg_engagement', title='Engagement Trends Over Time')
            st.plotly_chart(fig, use_container_width=True)

# ML Tabs (keeping existing ones)
with tabs[10]:
    st.header("ML Explainability (SHAP Summary Plot)")
    if st.button("Load SHAP Plot"):
        shap_url = "http://backend:8000/ml/explain/"
        st.image(shap_url)

with tabs[11]:
    st.header("Engagement Forecast (Next 30 Days)")
    response = requests.get("http://backend:8000/ml/forecast/")
    forecast_data = response.json() if response.status_code == 200 else []
    if forecast_data:
        df_forecast = pd.DataFrame(forecast_data)
        df_forecast['ds'] = pd.to_datetime(df_forecast['ds'])
        fig = px.line(df_forecast, x='ds', y='yhat', title='Forecasted Engagement')
        fig.add_scatter(x=df_forecast['ds'], y=df_forecast['yhat_upper'], mode='lines', line=dict(dash='dash'),
                        name='Upper Bound')
        fig.add_scatter(x=df_forecast['ds'], y=df_forecast['yhat_lower'], mode='lines', line=dict(dash='dash'),
                        name='Lower Bound')
        st.plotly_chart(fig)

with tabs[12]:
    st.header("Caption Clustering & Suggestions")
    response = requests.get("http://backend:8000/ml/cluster/")
    clusters = response.json() if response.status_code == 200 else {}
    for cluster_name, samples in clusters.items():
        st.subheader(cluster_name)
        for s in samples:
            st.write("-", s)

# Export Tab
with tabs[13]:
    st.header("Export Data & Professional Report")
    preview = requests.get("http://backend:8000/preview/")
    df = pd.DataFrame(preview.json())
    if len(df) > 0:
        csv = df.to_csv(index=False)
        st.download_button("Download CSV", csv, "social_media_analysis.csv", "text/csv")
        if st.button("Generate and Download PPTX Report"):
            tmp_report_path = "tmu_social_media_report.pptx"
            create_pptx_report(df, logo_path, tmp_report_path)
            with open(tmp_report_path, "rb") as f:
                st.download_button("Download TMU PPTX Report", f, "tmu_report.pptx")
